# from pptx import Presentation
# from pptx.util import Inches
# import os
# from datetime import datetime

# def generate_ppt(key_points: list[str], filename: str = "summary") -> str:
#     """
#     Generates a PowerPoint presentation with the given key points.
#     Each point is added as a bullet on a single slide.
#     Returns the path to the generated file.
#     """
#     # Create a presentation
#     prs = Presentation()

#     # Add a title slide
#     title_slide_layout = prs.slide_layouts[0]
#     slide = prs.slides.add_slide(title_slide_layout)
#     slide.shapes.title.text = "Meeting Summary"
#     slide.placeholders[1].text = f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M')}"

#     # Add a content slide with bullet points
#     bullet_slide_layout = prs.slide_layouts[1]
#     bullet_slide = prs.slides.add_slide(bullet_slide_layout)
#     shapes = bullet_slide.shapes
#     shapes.title.text = "Key Discussion Points"

#     body_shape = shapes.placeholders[1]
#     tf = body_shape.text_frame
#     tf.clear()

#     for point in key_points:
#         p = tf.add_paragraph()
#         p.text = point
#         p.level = 0

#     # Ensure output directory exists
#     output_dir = "static/presentations"
#     os.makedirs(output_dir, exist_ok=True)

#     # Save the presentation
#     output_path = os.path.join(output_dir, f"{filename}_summary.pptx")
#     prs.save(output_path)

#     return output_path

# from pptx import Presentation
# import os
# from datetime import datetime

# MAX_BULLETS_PER_SLIDE = 5  # Change this as needed

# def generate_ppt(key_points, filename="summary"):
#     prs = Presentation()

#     # Title Slide
#     slide = prs.slides.add_slide(prs.slide_layouts[0])
#     slide.shapes.title.text = "Meeting Summary"
#     slide.placeholders[1].text = f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M')}"

#     # Break key_points into chunks for multiple slides
#     for i in range(0, len(key_points), MAX_BULLETS_PER_SLIDE):
#         chunk = key_points[i:i + MAX_BULLETS_PER_SLIDE]

#         bullet_slide = prs.slides.add_slide(prs.slide_layouts[1])
#         bullet_slide.shapes.title.text = f"Key Discussion Points (Slide {i // MAX_BULLETS_PER_SLIDE + 1})"
        
#         tf = bullet_slide.shapes.placeholders[1].text_frame
#         tf.clear()

#         for point in chunk:
#             if point.strip():
#                 p = tf.add_paragraph()
#                 p.text = point.strip()
#                 p.level = 0

#     output_dir = "static/presentations"
#     os.makedirs(output_dir, exist_ok=True)
#     output_path = os.path.join(output_dir, f"{filename}_summary.pptx")
#     prs.save(output_path)

#     return output_path


from pptx import Presentation 
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from datetime import datetime
import os
import logging
import re

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

def parse_ai_ppt_output(ai_output):
    """Parse the AI-generated PPT output into a structured format."""
    sections = []
    current_title = None
    current_content = []

    lines = ai_output.split('\n')
    ppt_type = None
    num_slides = None
    i = 0

    # Extract PPT type and number of slides
    while i < len(lines):
        line = lines[i].strip()
        if not ppt_type and line.lower().startswith('**recommended ppt type**:'):
            i += 1
            while i < len(lines) and not lines[i].strip():
                i += 1
            if i < len(lines) and lines[i].strip():
                ppt_type = lines[i].strip()
                i += 1
        elif not num_slides and line.lower().startswith('**number of slides**:'):
            i += 1
            while i < len(lines) and not lines[i].strip():
                i += 1
            if i < len(lines) and lines[i].strip() and any(c.isdigit() for c in lines[i].strip()):
                num_slides_text = lines[i].strip()
                num_slides = int(''.join(filter(str.isdigit, num_slides_text)))
                i += 1
        elif line.lower().startswith('**slide breakdown**:'):
            i += 1
            break
        else:
            i += 1

    # Parse slide breakdown
    in_slide = False
    while i < len(lines):
        line = lines[i].strip()
        if line and line.startswith(('1. ', '2. ', '3. ', '4. ', '5. ', '6. ')) and '**' in line:
            if current_title and current_content:
                sections.append((current_title, '\n'.join(current_content)))
            current_title = line.split('**')[1].strip()
            current_content = []
            in_slide = True
        elif in_slide and line and not line.startswith(('1. ', '2. ', '3. ', '4. ', '5. ', '6. ')):
            current_content.append(line)
        elif not line and current_title and current_content:
            sections.append((current_title, '\n'.join(current_content)))
            current_content = []
        i += 1

    if current_title and current_content:
        sections.append((current_title, '\n'.join(current_content)))

    return ppt_type, num_slides, sections

def extract_title(text):
    for line in text.splitlines():
        if line.startswith("**Recommended PPT Type**:"):
            # Split on the first colon and take the second part
            return line.split(":", 1)[1].strip()
    return None

def set_background(slide, left, top, width, height):
    """Set a solid white background for the slide."""
    try:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 0, 0)
    except Exception as e:
        logging.error(f"Error setting background: {e}")

def add_cover_slide(prs, layout_idx, title, subtitle):
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        logging.info(f"Using layout index {layout_idx} for cover slide.")
    except IndexError:
        logging.warning(f"Layout index {layout_idx} not found. Using first available layout (index 0).")
        slide = prs.slides.add_slide(prs.slide_layouts[0])

    set_background(slide, Inches(0), Inches(0), Inches(10), Inches(7.5))

    title_shape = slide.shapes.add_textbox(Inches(0), Inches(1), Inches(8), Inches(2))
    tf = title_shape.text_frame
    p = tf.add_paragraph()
    p.text = "AutoReQ"
    p.font.size = Pt(72)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.CENTER
    logging.debug("Added title text box.")

    subtitle_shape = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(0.7))
    tf = subtitle_shape.text_frame
    p = tf.add_paragraph()
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.CENTER
    logging.debug("Added subtitle text box.")

    headline_shape = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.7))
    tf = headline_shape.text_frame
    p = tf.add_paragraph()
    p.text = "Meeting Transcript Analysis"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.CENTER
    logging.debug("Added headline text box.")

    date_shape = slide.shapes.add_textbox(Inches(6), Inches(6), Inches(3), Inches(0.5))
    tf = date_shape.text_frame
    p = tf.add_paragraph()
    p.text = datetime.now().strftime("Date %d-%m-%Y %I:%M %p IST")
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.RIGHT
    logging.debug("Added date text box.")

    footer_shape = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
    tf = footer_shape.text_frame
    p = tf.add_paragraph()
    p.text = "© Hitachi Digital Services 2025. All Rights Reserved"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.LEFT
    logging.debug("Added footer text box.")

    logging.info("Cover slide added successfully.")

def add_contents_slide(prs, layout_idx, sections):
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    except IndexError:
        logging.warning(f"Layout index {layout_idx} not found. Using first available layout.")
        slide = prs.slides.add_slide(prs.slide_layouts[0])

    set_background(slide, Inches(0), Inches(0), Inches(10), Inches(7.5))

    title_shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    tf = title_shape.text_frame
    p = tf.add_paragraph()
    p.text = "Agenda"
    p.font.size = Pt(36)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.CENTER

    content_shape = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    tf = content_shape.text_frame
    for i, (heading, _) in enumerate(sections, 1):
        p = tf.add_paragraph()
        p.text = f"{i}. {heading}"
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.LEFT

    footer_shape = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
    tf = footer_shape.text_frame
    p = tf.add_paragraph()
    p.text = "Hitachi Digital Services | Presentation headline title"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.LEFT

    copyright_shape = slide.shapes.add_textbox(Inches(6), Inches(6.5), Inches(3), Inches(0.5))
    tf = copyright_shape.text_frame
    p = tf.add_paragraph()
    p.text = "© Hitachi Digital Services 2025. All Rights Reserved"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.RIGHT

    logging.info("Contents slide added successfully.")

def add_content_slide(prs, layout_idx, heading, content):
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    except IndexError:
        logging.warning(f"Layout index {layout_idx} not found. Using first available layout.")
        slide = prs.slides.add_slide(prs.slide_layouts[0])

    set_background(slide, Inches(0), Inches(0), Inches(10), Inches(7.5))

    # Add title
    title_shape = slide.shapes.add_textbox(Inches(0), Inches(1), Inches(10), Inches(1))
    tf = title_shape.text_frame
    p = tf.add_paragraph()
    p.text = heading
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.CENTER

    # Add content with wrapping and multiple paragraphs
    content_height = Inches(4)  # Initial content height
    y_offset = Inches(2)  # Starting y-position for content
    content_lines = content.split('\n')
    max_paragraphs_per_box = 10  # Arbitrary limit to prevent overflow
    current_paragraph_count = 0
    current_box = slide.shapes.add_textbox(Inches(1), y_offset, Inches(8), content_height)
    tf = current_box.text_frame
    tf.word_wrap = True  # Enable word wrapping

    for line in content_lines:
        if line.strip():
            if current_paragraph_count >= max_paragraphs_per_box:
                y_offset += content_height
                current_box = slide.shapes.add_textbox(Inches(1), y_offset, Inches(8), content_height)
                tf = current_box.text_frame
                tf.word_wrap = True
                current_paragraph_count = 0
            p = tf.add_paragraph()
            p.text = f"• {line.strip()}"
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.alignment = PP_ALIGN.LEFT
            current_paragraph_count += 1

    # Add footer
    footer_shape = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
    tf = footer_shape.text_frame
    p = tf.add_paragraph()
    p.text = "Hitachi Digital Services | Presentation headline title"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.LEFT

    copyright_shape = slide.shapes.add_textbox(Inches(6), Inches(6.5), Inches(3), Inches(0.5))
    tf = copyright_shape.text_frame
    p = tf.add_paragraph()
    p.text = "© Hitachi Digital Services 2025. All Rights Reserved"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.RIGHT

    logging.info(f"Content slide '{heading}' added successfully.")

def add_closing_slide(prs, layout_idx):
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    except IndexError:
        logging.warning(f"Layout index {layout_idx} not found. Using first available layout.")
        slide = prs.slides.add_slide(prs.slide_layouts[0])

    set_background(slide, Inches(0), Inches(0), Inches(10), Inches(7.5))

    title_shape = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
    tf = title_shape.text_frame
    p = tf.add_paragraph()
    p.text = "Thank You"
    p.font.size = Pt(40)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.CENTER

    footer_shape = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
    tf = footer_shape.text_frame
    p = tf.add_paragraph()
    p.text = "Hitachi Digital Services | Presentation headline title"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.LEFT

    copyright_shape = slide.shapes.add_textbox(Inches(6), Inches(6.5), Inches(3), Inches(0.5))
    tf = copyright_shape.text_frame
    p = tf.add_paragraph()
    p.text = "© Hitachi Digital Services 2025. All Rights Reserved"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.RIGHT

    logging.info("Closing slide added successfully.")

def create_presentation(ai_output, template_path, output_path, title, subtitle):
    """Main function to create the PowerPoint presentation from AI output."""
    try:
        if not os.path.exists(template_path):
            logging.warning(f"Template file not found at: {template_path}. Creating new presentation.")
            prs = Presentation()
        else:
            prs = Presentation(template_path)
    except Exception as e:
        logging.error(f"Error loading template: {e}. Creating new presentation.")
        prs = Presentation()

    ppt_type, num_slides, sections = parse_ai_ppt_output(ai_output)
    if not sections:
        logging.error("No sections parsed from AI output. Presentation may be empty.")
        return

    logging.debug(f"Parsed sections: {sections}")  # Debug parsed sections
    blank_layout_idx = len(prs.slide_layouts) -6
    if blank_layout_idx < 0:
        logging.warning("No valid layouts found. Using default layout (index 0).")
        blank_layout_idx = 0
    logging.info(f"Using layout index {blank_layout_idx} for all slides.")

    try:
        add_cover_slide(prs, blank_layout_idx, title, subtitle)
        add_contents_slide(prs, blank_layout_idx, sections)
        for heading, content in sections[:num_slides]:  # Limit to number of slides specified
            logging.debug(f"Adding slide: {heading}, Content: {content}")
            add_content_slide(prs, blank_layout_idx, heading, content)
        add_closing_slide(prs, blank_layout_idx)
    except Exception as e:
        logging.error(f"Error adding slides: {e}")
        return

    try:
        prs.save(output_path)
        logging.info(f"Presentation saved to {output_path}")
    except Exception as e:
        logging.error(f"Error saving presentation: {e}")

